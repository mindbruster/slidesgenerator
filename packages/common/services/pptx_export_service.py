"""
PPTX Export Service - generates editable PowerPoint presentations
"""

from io import BytesIO
from typing import TypedDict

import httpx
from pptx import Presentation as PptxPresentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from packages.common.models import Presentation, Slide
from packages.common.themes import get_theme


class ThemeColors(TypedDict):
    """Colors used for PPTX export"""

    bg: RGBColor
    accent: RGBColor
    text: RGBColor
    muted: RGBColor


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor"""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Content margins
MARGIN_LEFT = Inches(1)
MARGIN_TOP = Inches(1)
CONTENT_WIDTH = Inches(11.333)
CONTENT_HEIGHT = Inches(5.5)


def generate_pptx(presentation: Presentation) -> bytes:
    """
    Generate a PPTX file from a presentation.

    Args:
        presentation: The presentation model with slides

    Returns:
        bytes: The PPTX file as bytes
    """
    # Get theme colors
    theme = get_theme(presentation.theme or "neobrutalism")
    colors: ThemeColors = {
        "bg": _hex_to_rgb(theme.colors.background),
        "accent": _hex_to_rgb(theme.colors.accent),
        "text": _hex_to_rgb(theme.colors.text_primary),
        "muted": _hex_to_rgb(theme.colors.text_secondary),
    }

    prs = PptxPresentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Sort slides by order
    sorted_slides = sorted(presentation.slides, key=lambda s: s.order)

    for slide in sorted_slides:
        _add_slide(prs, slide, colors)

    # Write to bytes
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()


def _add_slide(prs: PptxPresentation, slide: Slide, colors: ThemeColors) -> None:
    """Add a slide to the presentation based on its type"""
    if slide.type == "title":
        _add_title_slide(prs, slide, colors)
    elif slide.type == "bullets":
        _add_bullets_slide(prs, slide, colors)
    elif slide.type == "quote":
        _add_quote_slide(prs, slide, colors)
    elif slide.type == "section":
        _add_section_slide(prs, slide, colors)
    elif slide.type == "chart":
        _add_chart_slide(prs, slide, colors)
    else:
        _add_content_slide(prs, slide, colors)


def _add_title_slide(prs: PptxPresentation, slide: Slide, colors: ThemeColors) -> None:
    """Add a title slide with centered title and subtitle"""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    pptx_slide = prs.slides.add_slide(blank_layout)
    _set_slide_background(pptx_slide, colors)

    # Title - centered vertically and horizontally
    if slide.title:
        title_box = pptx_slide.shapes.add_textbox(
            MARGIN_LEFT,
            Inches(2.5),
            CONTENT_WIDTH,
            Inches(1.5),
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = slide.title
        title_para.font.size = Pt(60)
        title_para.font.bold = True
        title_para.font.color.rgb = colors["text"]
        title_para.alignment = PP_ALIGN.CENTER
        title_frame.auto_size = None

    # Subtitle
    if slide.subtitle:
        sub_box = pptx_slide.shapes.add_textbox(
            MARGIN_LEFT,
            Inches(4.2),
            CONTENT_WIDTH,
            Inches(1),
        )
        sub_frame = sub_box.text_frame
        sub_frame.word_wrap = True
        sub_para = sub_frame.paragraphs[0]
        sub_para.text = slide.subtitle
        sub_para.font.size = Pt(28)
        sub_para.font.color.rgb = colors["muted"]
        sub_para.alignment = PP_ALIGN.CENTER


def _add_slide_image(pptx_slide, slide: Slide, left, top, width, height) -> None:
    """Add an image to a slide by downloading from URL"""
    if not slide.image_url:
        return

    try:
        # Download image
        with httpx.Client() as client:
            response = client.get(slide.image_url, timeout=10.0)
            response.raise_for_status()
            image_data = BytesIO(response.content)

        # Add image to slide
        pptx_slide.shapes.add_picture(image_data, left, top, width, height)
    except Exception as e:
        # Log error but don't fail the export
        print(f"Failed to add image to slide: {e}")


def _add_bullets_slide(prs: PptxPresentation, slide: Slide, colors: ThemeColors) -> None:
    """Add a slide with title and bullet points"""
    blank_layout = prs.slide_layouts[6]
    pptx_slide = prs.slides.add_slide(blank_layout)
    _set_slide_background(pptx_slide, colors)

    has_image = bool(slide.image_url)
    text_width = Inches(5.5) if has_image else CONTENT_WIDTH

    # Title at top
    if slide.title:
        title_box = pptx_slide.shapes.add_textbox(
            MARGIN_LEFT,
            MARGIN_TOP,
            text_width,
            Inches(1),
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide.title
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = colors["text"]

    # Bullets
    if slide.bullets:
        bullets_box = pptx_slide.shapes.add_textbox(
            MARGIN_LEFT,
            Inches(2.2),
            text_width,
            Inches(4.5),
        )
        bullets_frame = bullets_box.text_frame
        bullets_frame.word_wrap = True

        for i, bullet in enumerate(slide.bullets):
            if i == 0:
                para = bullets_frame.paragraphs[0]
            else:
                para = bullets_frame.add_paragraph()
            para.text = f"• {bullet}"
            para.font.size = Pt(24)
            para.font.color.rgb = colors["text"]
            para.space_after = Pt(16)

    # Add image if present
    if has_image:
        _add_slide_image(pptx_slide, slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5))


def _add_content_slide(prs: PptxPresentation, slide: Slide, colors: ThemeColors) -> None:
    """Add a slide with title and body text"""
    blank_layout = prs.slide_layouts[6]
    pptx_slide = prs.slides.add_slide(blank_layout)
    _set_slide_background(pptx_slide, colors)

    has_image = bool(slide.image_url)
    text_width = Inches(5.5) if has_image else CONTENT_WIDTH

    # Title
    if slide.title:
        title_box = pptx_slide.shapes.add_textbox(
            MARGIN_LEFT,
            MARGIN_TOP,
            text_width,
            Inches(1),
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide.title
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = colors["text"]

    # Body text
    if slide.body:
        body_box = pptx_slide.shapes.add_textbox(
            MARGIN_LEFT,
            Inches(2.2),
            text_width,
            Inches(4.5),
        )
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        body_para = body_frame.paragraphs[0]
        body_para.text = slide.body
        body_para.font.size = Pt(24)
        body_para.font.color.rgb = colors["text"]
        body_para.line_spacing = 1.5

    # Add image if present
    if has_image:
        _add_slide_image(pptx_slide, slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5))


def _add_quote_slide(prs: PptxPresentation, slide: Slide, colors: ThemeColors) -> None:
    """Add a slide with a styled quote and attribution"""
    blank_layout = prs.slide_layouts[6]
    pptx_slide = prs.slides.add_slide(blank_layout)
    _set_slide_background(pptx_slide, colors)

    # Accent bar on left
    accent_bar = pptx_slide.shapes.add_shape(
        1,  # Rectangle
        Inches(1),
        Inches(2.5),
        Inches(0.08),
        Inches(2.5),
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = colors["accent"]
    accent_bar.line.fill.background()

    # Quote text
    if slide.quote:
        quote_box = pptx_slide.shapes.add_textbox(
            Inches(1.4),
            Inches(2.5),
            Inches(10),
            Inches(2),
        )
        quote_frame = quote_box.text_frame
        quote_frame.word_wrap = True
        quote_para = quote_frame.paragraphs[0]
        quote_para.text = f'"{slide.quote}"'
        quote_para.font.size = Pt(32)
        quote_para.font.italic = True
        quote_para.font.color.rgb = colors["text"]

    # Attribution
    if slide.attribution:
        attr_box = pptx_slide.shapes.add_textbox(
            Inches(1.4),
            Inches(5),
            Inches(10),
            Inches(0.5),
        )
        attr_frame = attr_box.text_frame
        attr_para = attr_frame.paragraphs[0]
        attr_para.text = f"— {slide.attribution}"
        attr_para.font.size = Pt(20)
        attr_para.font.color.rgb = colors["muted"]


def _add_section_slide(prs: PptxPresentation, slide: Slide, colors: ThemeColors) -> None:
    """Add a section divider slide with centered text"""
    blank_layout = prs.slide_layouts[6]
    pptx_slide = prs.slides.add_slide(blank_layout)
    _set_slide_background(pptx_slide, colors)

    if slide.title:
        title_box = pptx_slide.shapes.add_textbox(
            MARGIN_LEFT,
            Inches(3),
            CONTENT_WIDTH,
            Inches(1.5),
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = slide.title
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = colors["text"]
        title_para.alignment = PP_ALIGN.CENTER


def _add_chart_slide(prs: PptxPresentation, slide: Slide, colors: ThemeColors) -> None:
    """Add a slide with a professionally styled chart"""
    blank_layout = prs.slide_layouts[6]
    pptx_slide = prs.slides.add_slide(blank_layout)
    _set_slide_background(pptx_slide, colors)

    # Title
    if slide.title:
        title_box = pptx_slide.shapes.add_textbox(
            MARGIN_LEFT,
            MARGIN_TOP,
            CONTENT_WIDTH,
            Inches(1),
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide.title
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = colors["text"]

    # Chart
    if slide.chart_data and slide.chart_type:
        chart_data = CategoryChartData()

        # Extract labels and values from chart_data
        labels = [point.get("label", "") for point in slide.chart_data]
        values = [point.get("value", 0) for point in slide.chart_data]

        chart_data.categories = labels
        chart_data.add_series("Data", values)

        # Map chart type to pptx chart type
        chart_type_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "horizontal_bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,  # Use markers for better visibility
            "pie": XL_CHART_TYPE.PIE,
            "donut": XL_CHART_TYPE.DOUGHNUT,
            "area": XL_CHART_TYPE.AREA,
        }
        xl_chart_type = chart_type_map.get(slide.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

        # Add chart with better positioning
        chart_left = Inches(1.5)
        chart_top = Inches(2)
        chart_width = Inches(10)
        chart_height = Inches(5)

        chart_shape = pptx_slide.shapes.add_chart(
            xl_chart_type,
            chart_left,
            chart_top,
            chart_width,
            chart_height,
            chart_data,
        )

        # Style the chart
        chart = chart_shape.chart

        # Style the plot area
        plot = chart.plots[0]

        # Add data labels for better readability
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(11)
        data_labels.font.bold = True
        data_labels.font.color.rgb = colors["text"]

        # Position data labels based on chart type
        if slide.chart_type in ("pie", "donut"):
            data_labels.show_percentage = True
            data_labels.show_value = False
            data_labels.show_category_name = True
        else:
            data_labels.show_value = True

        # Style the series with theme colors
        if len(chart.series) > 0:
            series = chart.series[0]

            # Set fill color for bars/areas
            if slide.chart_type in ("bar", "horizontal_bar", "area"):
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = colors["accent"]

            # Set line color for line charts
            if slide.chart_type == "line":
                series.format.line.color.rgb = colors["accent"]
                series.format.line.width = Pt(3)
                # Style markers
                series.marker.style = 2  # Circle marker
                series.marker.size = 10
                series.marker.format.fill.solid()
                series.marker.format.fill.fore_color.rgb = colors["accent"]
                series.marker.format.line.color.rgb = colors["text"]

            # Apply individual colors to pie/donut slices if provided
            if slide.chart_type in ("pie", "donut"):
                # Generate a nice color palette
                base_colors = [
                    colors["accent"],
                    _adjust_color(colors["accent"], 0.7),
                    _adjust_color(colors["accent"], 0.5),
                    _adjust_color(colors["accent"], 0.3),
                    colors["muted"],
                ]
                for i, point in enumerate(series.points):
                    point.format.fill.solid()
                    color_idx = i % len(base_colors)
                    point.format.fill.fore_color.rgb = base_colors[color_idx]

        # Style category axis (x-axis)
        if hasattr(chart, 'category_axis') and chart.category_axis:
            cat_axis = chart.category_axis
            cat_axis.tick_labels.font.size = Pt(11)
            cat_axis.tick_labels.font.color.rgb = colors["text"]
            if hasattr(cat_axis, 'format') and hasattr(cat_axis.format, 'line'):
                cat_axis.format.line.color.rgb = colors["muted"]

        # Style value axis (y-axis)
        if hasattr(chart, 'value_axis') and chart.value_axis:
            val_axis = chart.value_axis
            val_axis.tick_labels.font.size = Pt(11)
            val_axis.tick_labels.font.color.rgb = colors["text"]
            if hasattr(val_axis, 'format') and hasattr(val_axis.format, 'line'):
                val_axis.format.line.color.rgb = colors["muted"]
            # Add gridlines
            val_axis.has_major_gridlines = True
            if val_axis.major_gridlines:
                val_axis.major_gridlines.format.line.color.rgb = _hex_to_rgb("#e0e0e0")
                val_axis.major_gridlines.format.line.width = Pt(0.5)

        # Hide legend for cleaner look (data labels show values)
        chart.has_legend = False

    else:
        # Fallback: show message if no chart data
        msg_box = pptx_slide.shapes.add_textbox(
            MARGIN_LEFT,
            Inches(3.5),
            CONTENT_WIDTH,
            Inches(1),
        )
        msg_frame = msg_box.text_frame
        msg_para = msg_frame.paragraphs[0]
        msg_para.text = "[Chart data not available]"
        msg_para.font.size = Pt(20)
        msg_para.font.color.rgb = colors["muted"]
        msg_para.alignment = PP_ALIGN.CENTER


def _adjust_color(base_color: RGBColor, factor: float) -> RGBColor:
    """Adjust color brightness for creating color variations"""
    r = min(255, int(base_color[0] + (255 - base_color[0]) * (1 - factor)))
    g = min(255, int(base_color[1] + (255 - base_color[1]) * (1 - factor)))
    b = min(255, int(base_color[2] + (255 - base_color[2]) * (1 - factor)))
    return RGBColor(r, g, b)


def _set_slide_background(pptx_slide, colors: ThemeColors) -> None:
    """Set the background color for a slide based on theme"""
    background = pptx_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors["bg"]
